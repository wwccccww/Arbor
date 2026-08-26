from __future__ import annotations

import io

from arbor.adapters.outbound.multimodal.types import MediaChunk, MediaParseResult


def parse_pptx(data: bytes, filename: str) -> MediaParseResult:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    chunks: list[MediaChunk] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text.strip())
        text = "\n".join(lines).strip()
        if not text:
            continue
        chunks.append(
            MediaChunk(
                text=text,
                memory_type="file_chunk",
                metadata={"source": filename, "slide": slide_no, "parser": "python-pptx"},
            )
        )
    return MediaParseResult(chunks=chunks, parser="python-pptx", media_kind="document")
